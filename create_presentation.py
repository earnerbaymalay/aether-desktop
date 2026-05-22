from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_slide(prs, title, content, image_path=None):
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background to black
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 10)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 215, 0) # Gold/Yellow
    
    # Add content
    if image_path:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
        slide.shapes.add_picture(image_path, Inches(5), Inches(2), height=Inches(4.5))
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for bullet in content:
        p = content_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(240, 240, 240) # Off-white
        p.space_after = Pt(14)
        p.level = 0
        
def create_presentation():
    prs = Presentation()
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title Slide
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 10)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "NODA"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 215, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "THE DETACHABLE SOVEREIGNTY HUB\n(NOCTA STANDARDS)"
    p2.font.size = Pt(32)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER

    # Vision & Mission
    create_slide(prs, "VISION & MISSION", [
        "Modular Sovereignty: A high-end, privacy-first alternative to Google Home and Amazon Echo.",
        "Zero Telemetry by Default: Designed to appeal to GDPR/HIPAA-adjacent requirements.",
        "Acoustic Silence: Passive cooling derived from structural heatsink engineering.",
        "Vendor Agnostic: Hardware independent, built on a hardened Linux ecosystem.",
        "Nocta Aesthetic: Premium matte black aluminum with bright yellow accents and glassmorphism."
    ])

    # Hardware Concept
    create_slide(prs, "THE MONOLITH: HARDWARE CONCEPT", [
        "Chassis: Sandblasted Matte Obsidian aluminum.",
        "Heatsink: Blackened copper fins for silent, passive cooling.",
        "Display: Detachable 10-inch glassmorphic OLED screen.",
        "Accents: Gold/Yellow LED status indicators.",
        "Privacy: Physical sliding toggles and kill-switches for microphones and cameras."
    ], "/Users/brookeconnor/.gemini/antigravity/brain/9543bdfd-3b6a-4ad2-85fa-ac2dbc529415/noda_hardware_concept_1779144228699.png")

    # Software & UI
    create_slide(prs, "SOFTWARE: THE FORGE & VIBE-CODED UI", [
        "The Immutable Core: Built on Alpine Linux with a read-only root filesystem.",
        "Generative Sandbox ('The Forge'): A vibe-coded engine that builds apps via voice/text.",
        "Neural Sovereignty: 100% local LLM execution with interchangeable model brains (Llama, Qwen, Mistral).",
        "Interface: Glassmorphism with terminal-inspired typography (Nocta Dark Mode)."
    ], "/Users/brookeconnor/.gemini/antigravity/brain/9543bdfd-3b6a-4ad2-85fa-ac2dbc529415/noda_software_ui_1779144291289.png")

    # Expansion Ecosystem
    create_slide(prs, "EXPANSION ECOSYSTEM: MODULAR BRICKS", [
        "Hot-swappable hardware modules that stack or dock to the core via magnetic gold pogo pins.",
        "SDR / Radio Brick: Adds Zigbee, Z-Wave, and localized RF control for secure networking.",
        "NVMe Storage Brick: High-capacity local data vault (converts the core into an encrypted NAS).",
        "Audiophile DAC Brick: High-fidelity audio processing for localized media."
    ], "/Users/brookeconnor/.gemini/antigravity/brain/9543bdfd-3b6a-4ad2-85fa-ac2dbc529415/noda_expansion_bricks_1779144315748.png")

    # Market Analysis
    create_slide(prs, "MARKET ANALYSIS & TARGET AUDIENCE", [
        "The Sovereign Prosumer: A growing demographic of privacy-focused, 'De-Googled' users.",
        "Homelab Enthusiasts: Users who find existing hubs too restrictive for advanced local automation.",
        "B2B Licensing: Providing a secure, private hub for corporate offices or high-security environments.",
        "Competitive Advantage: Extensibility, offline capability, and 100% data ownership."
    ])

    # Production Rollout
    create_slide(prs, "PRODUCTION ROLLOUT ROADMAP", [
        "Q3 2026: 'Headless Founders' Beta (No Screen) release for developers and stress testing.",
        "Q1 2027: Limited 'Chromal Stealth / Obsidian' Complete Bundle Launch (Exclusivity-driven).",
        "Q3 2027: Rollout of NVMe and SDR Radio Expansion Bricks.",
        "Q1 2028: Sovereign App Marketplace Launch (Vetted FOSS, local-only AI apps)."
    ])

    prs.save("/Users/brookeconnor/Downloads/drive-download-20260518T215206Z-3-001/NODA_Nocta_Presentation.pptx")
    print("Presentation saved successfully at /Users/brookeconnor/Downloads/drive-download-20260518T215206Z-3-001/NODA_Nocta_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
